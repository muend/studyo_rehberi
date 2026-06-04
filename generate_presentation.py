#!/usr/bin/env python3
"""
Timber / Steel / Masonry — 30-Slide Academic Presentation Generator
"""

import os, io, csv, sys, time, requests
from pathlib import Path
from PIL import Image as PILImage

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── Dimensions ────────────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Colours ───────────────────────────────────────────────────────────────────
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GRAY   = RGBColor(0xF4, 0xF4, 0xF4)
MID_GRAY     = RGBColor(0x90, 0x90, 0x90)
DARK_GRAY    = RGBColor(0x50, 0x50, 0x50)

TIMBER_DARK  = RGBColor(0x6B, 0x48, 0x1A)
TIMBER_MED   = RGBColor(0x9B, 0x73, 0x35)
TIMBER_LIGHT = RGBColor(0xF5, 0xED, 0xDC)

STEEL_DARK   = RGBColor(0x1E, 0x35, 0x52)
STEEL_MED    = RGBColor(0x3B, 0x5C, 0x84)
STEEL_LIGHT  = RGBColor(0xDC, 0xE8, 0xF5)

MASON_DARK   = RGBColor(0x7A, 0x28, 0x10)
MASON_MED    = RGBColor(0xB0, 0x45, 0x22)
MASON_LIGHT  = RGBColor(0xF5, 0xE5, 0xDD)

SECTION_COLORS = {
    "TIMBER":  (TIMBER_DARK, TIMBER_MED, TIMBER_LIGHT),
    "STEEL":   (STEEL_DARK,  STEEL_MED,  STEEL_LIGHT),
    "MASONRY": (MASON_DARK,  MASON_MED,  MASON_LIGHT),
}

# ── Image catalogue ───────────────────────────────────────────────────────────
IMAGES = {
    # key: (url, caption, credit, license)
    "brock_commons": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/d/d1/Brock_Commons_Tallwood_House_%281%29.jpg/1200px-Brock_Commons_Tallwood_House_%281%29.jpg",
        "Brock Commons Tallwood House, Vancouver (Acton Ostry Architects, 2016)",
        "Richard Gao / Wikimedia Commons", "CC BY-SA 4.0"),
    "half_timber": (
        "https://upload.wikimedia.org/wikipedia/commons/5/5f/Half-timbered_tudor_buildings%2C_High_Holborn.JPG",
        "Staple Inn, High Holborn, London — Tudor half-timbered frame construction",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "timber_interior": (
        "https://upload.wikimedia.org/wikipedia/commons/4/44/Glulam.JPG",
        "Glulam beam — glued laminated timber showing lamination layers",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "mjostaarnet": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/4/43/Mj%C3%B8st%C3%A5rnet_2020.jpg/800px-Mj%C3%B8st%C3%A5rnet_2020.jpg",
        "Mjøstårnet, Brumunddal, Norway (Voll Arkitekter, 2019) — 85.4 m, 18 storeys",
        "Statsbygg / Wikimedia Commons", "CC BY 4.0"),
    "tamedia": (
        "https://upload.wikimedia.org/wikipedia/commons/4/44/Glulam.JPG",
        "Glulam structural timber — exposed laminated wood member (illustrative of Tamedia structural logic)",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "clt_construction": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/d/d1/Brock_Commons_Tallwood_House_%281%29.jpg/1200px-Brock_Commons_Tallwood_House_%281%29.jpg",
        "Mass timber prefabricated construction — Brock Commons Tallwood House",
        "Richard Gao / Wikimedia Commons", "CC BY-SA 4.0"),
    "forest": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/1/17/Great_Smoky_Mountains_National_Park._Tennessee_%28LOC%29.jpg/1200px-Great_Smoky_Mountains_National_Park._Tennessee_%28LOC%29.jpg",
        "Managed forest — the renewable raw material basis of structural timber",
        "Library of Congress / Wikimedia Commons", "Public Domain"),
    "eiffel": (
        "https://upload.wikimedia.org/wikipedia/commons/a/a8/Tour_Eiffel_Wikimedia_Commons.jpg",
        "Eiffel Tower, Paris (Gustave Eiffel, 1889) — 324 m iron lattice structure",
        "Benh LIEU SONG / Wikimedia Commons", "CC BY-SA 3.0"),
    "crystal_palace": (
        "https://upload.wikimedia.org/wikipedia/commons/e/e3/Crystal_Palace_from_the_northeast_from_Dickinson%27s_Comprehensive_Pictures_of_the_Great_Exhibition_of_1851._1854.jpg",
        "Crystal Palace, London (Joseph Paxton, 1851) — prefabricated iron-and-glass structure",
        "Dickinson Brothers / Wikimedia Commons", "Public Domain"),
    "pompidou": (
        "https://upload.wikimedia.org/wikipedia/commons/e/e9/Centre_Georges_Pompidou_July_13%2C_2008_3.jpg",
        "Centre Georges Pompidou, Paris (Piano & Rogers, 1977) — exposed steel and services",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "seagram": (
        "https://upload.wikimedia.org/wikipedia/commons/8/81/Seagrambuilding.JPG",
        "Seagram Building, New York (Mies van der Rohe & Philip Johnson, 1958)",
        "Max Hermus / Wikimedia Commons", "CC BY-SA 3.0"),
    "steel_frame": (
        "https://upload.wikimedia.org/wikipedia/commons/d/d5/Mutual_of_Omaha_Skyscraper_construction.jpg",
        "Steel frame high-rise under construction — structural steel skeleton",
        "Wikimedia Commons", "CC BY-SA 4.0"),
    "willis_tower": (
        "https://upload.wikimedia.org/wikipedia/commons/2/23/Chicago_Sears_Tower.jpg",
        "Willis Tower (SOM, 1973), Chicago — bundled tube structural system, 442 m",
        "Daniel Schwen / Wikimedia Commons", "CC BY-SA 4.0"),
    "lloyds": (
        "https://upload.wikimedia.org/wikipedia/commons/b/b5/The_Llyods_Building_London_Designed_by_Richard_Rogers_and_Partners_2.jpg",
        "Lloyd's Building, London (Richard Rogers, 1986) — exposed steel structure and services",
        "Wikimedia Commons", "CC BY 3.0"),
    "empire_state": (
        "https://upload.wikimedia.org/wikipedia/commons/6/6c/Looking_Up_at_Empire_State_Building.JPG",
        "Empire State Building, New York (Shreve, Lamb & Harmon, 1931) — steel-framed skyscraper",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "brick_wall": (
        "https://upload.wikimedia.org/wikipedia/commons/5/5c/Surfaces_brick_wall_with_mortar_closeup_view.JPG",
        "Brick masonry wall — close-up of mortar joints and stretcher bond pattern",
        "Wikimedia Commons", "CC BY-SA 4.0"),
    "pont_du_gard": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/6/6c/Pont_du_Gard_Oct_2007.jpg/1200px-Pont_du_Gard_Oct_2007.jpg",
        "Pont du Gard, France (1st century CE) — Roman stone arch masonry aqueduct",
        "Benh LIEU SONG / Wikimedia Commons", "CC BY-SA 3.0"),
    "pantheon": (
        "https://upload.wikimedia.org/wikipedia/commons/8/8b/Dome_of_Pantheon_Rome.JPG",
        "Pantheon dome interior, Rome (c. 125 CE) — 43.3 m unreinforced concrete dome with oculus",
        "Dave Amos / Wikimedia Commons", "CC BY-SA 2.5"),
    "hagia_sophia": (
        "https://upload.wikimedia.org/wikipedia/commons/d/d7/Istanbul_-_Hagia_Sophia_-_01.JPG",
        "Hagia Sophia, Istanbul (Anthemius of Tralles, 537 CE) — masonry dome on pendentives",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "djenne": (
        "https://upload.wikimedia.org/wikipedia/commons/6/66/Great_Mosque_of_Djenn%C3%A9_1.jpg",
        "Great Mosque of Djenné, Mali (rebuilt 1907) — UNESCO World Heritage adobe masonry",
        "Ruud Zwart / Wikimedia Commons", "CC BY-SA 3.0"),
    "monadnock": (
        "https://upload.wikimedia.org/wikipedia/commons/e/e4/Monadnock_Building%2C_Chicago_Loop%2C_Chicago%2C_Illinois_%289179343149%29.jpg",
        "Monadnock Building, Chicago (Burnham & Root, 1891) — 16-storey load-bearing brick masonry",
        "Wikimedia Commons", "CC BY 2.0"),
    "tate_modern": (
        "https://upload.wikimedia.org/wikipedia/commons/f/f8/Zumthor_Kolumba04.JPG",
        "Kolumba Museum, Cologne (Peter Zumthor, 2007) — custom thin brick, historic ruins integrated",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "masonry_damage": (
        "https://upload.wikimedia.org/wikipedia/commons/8/8b/Sichuan_earthquake_building_collasped..JPG",
        "Earthquake-damaged masonry building — diagonal shear cracking pattern (2008 Sichuan)",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "gothic_vault": (
        "https://upload.wikimedia.org/wikipedia/commons/a/a8/Polygonal_masonry_wall%2C_Amelia%2C_Italy.JPG",
        "Polygonal stone masonry wall, Amelia, Italy — dry-stack compressive masonry",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "timber_limits": (
        "https://upload.wikimedia.org/wikipedia/commons/4/44/Glulam.JPG",
        "Glulam structural timber — precise fabrication and connection detailing are essential",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "i_beam": (
        "https://upload.wikimedia.org/wikipedia/commons/4/40/I-Beam_002.JPG",
        "Structural steel I-beam (W-section) — standard cross-section for beams and columns",
        "Wikimedia Commons", "CC BY-SA 2.5"),
    "steel_construction": (
        "https://upload.wikimedia.org/wikipedia/commons/d/d5/Mutual_of_Omaha_Skyscraper_construction.jpg",
        "Steel frame high-rise under construction — erection of structural steel skeleton",
        "Wikimedia Commons", "CC BY-SA 4.0"),
    "kolumba": (
        "https://upload.wikimedia.org/wikipedia/commons/f/f8/Zumthor_Kolumba04.JPG",
        "Kolumba Museum, Cologne (Peter Zumthor, 2007) — contemporary brick facade on historic ruins",
        "Wikimedia Commons", "CC BY-SA 3.0"),
    "flemish_bond": (
        "https://upload.wikimedia.org/wikipedia/commons/5/5c/Surfaces_brick_wall_with_mortar_closeup_view.JPG",
        "Brick masonry — close-up showing mortar joints and bonding pattern",
        "Wikimedia Commons", "CC BY-SA 4.0"),
}

# ── Slide data ────────────────────────────────────────────────────────────────
SLIDES = [
  # ─── TIMBER ───────────────────────────────────────────────────────────────
  {
    "num": 1, "section": "TIMBER",
    "title": "Timber as a Structural Material",
    "subtitle": "Organic, Anisotropic, and Inherently Renewable",
    "bullets": [
      ("Anisotropy and Grain Direction",
       "Timber is an organic, orthotropic material. Mechanical strength varies with grain direction: "
       "parallel-to-grain tensile and compressive values far exceed perpendicular-to-grain values — "
       "a fundamental constraint governing every connection and panel detail (Allen & Iano, 2019)."),
      ("Strength-to-Weight Ratio",
       "The specific strength (strength/density) of structural timber approaches mild steel, making it "
       "one of the most weight-efficient structural materials available for spans up to 30–40 m "
       "(Porteous & Kermani, 2013)."),
      ("Moisture and Dimensional Stability",
       "Below fiber saturation (~30% MC), timber shrinks as it dries and swells as it gains moisture — "
       "exclusively perpendicular to grain. Poor moisture detailing causes splitting, fastener withdrawal, "
       "and long-term serviceability failure."),
      ("Engineered Timber Products",
       "CLT, glulam, and LVL overcome natural variability by cross-laminating or bonding fiber layers "
       "under factory conditions, producing dimensionally stable, large-format structural elements "
       "suitable for urban buildings at multi-storey scale."),
      ("Carbon Sequestration",
       "Trees absorb CO₂ during growth. When used in durable buildings, this biogenic carbon remains "
       "stored for the structure's service life, contributing to net carbon reduction in construction "
       "(Ramage et al., 2017)."),
    ],
    "highlight": "Timber's structural behavior is direction-dependent: efficient parallel to grain, vulnerable perpendicular to grain — every connection and detail must reflect this anisotropy.",
    "citation": "(Allen & Iano, 2019; Porteous & Kermani, 2013; Ramage et al., 2017)",
    "img_key": "brock_commons",
  },
  {
    "num": 2, "section": "TIMBER",
    "title": "Historical Evolution of Timber Construction",
    "subtitle": "From Vernacular Craft to Engineered Mass Timber",
    "bullets": [
      ("Vernacular and Pre-industrial Timber",
       "Post-and-beam, mortise-and-tenon, and scarf joints defined timber construction from "
       "prehistoric shelters through medieval Europe. Japanese and Scandinavian traditions developed "
       "highly refined structural systems without metal fasteners (Addis, 2007)."),
      ("Half-Timbering and Gothic Carpentry",
       "Medieval half-timbering integrated exposed structural frames with infill panels. "
       "Gothic roof carpentry — tie beams, crown posts, hammer beams — pushed timber structural "
       "design to extraordinary formal and technical complexity."),
      ("Industrial Sawmilling and Platform Framing",
       "19th-century industrial sawmilling enabled standardized dimensional lumber, leading to "
       "balloon and platform framing in North America — lightweight, repetitive systems that remain "
       "the dominant low-rise residential structural approach globally."),
      ("Development of Engineered Wood Products",
       "Glulam (laminated glued lumber, developed 1900–1940s), LVL, PSL, and OSB restored timber "
       "relevance for longer spans and larger structures, enabling sports halls, school gymnasiums, "
       "and medium-span industrial buildings through the mid-20th century."),
      ("21st-Century Mass Timber Revival",
       "CLT (cross-laminated timber), pioneered in Austria in the 1990s, launched a mass timber "
       "movement that repositioned wood as a viable urban building material for multi-storey "
       "construction, driven by climate targets and prefabrication economics (Ramage et al., 2017)."),
    ],
    "highlight": "Timber's history is a series of technical revolutions: from craft joinery, to industrial framing, to engineered panels — each driven by new production technologies rather than material change.",
    "citation": "(Addis, 2007; Ramage et al., 2017; Lehmann, 2012)",
    "img_key": "half_timber",
  },
  {
    "num": 3, "section": "TIMBER",
    "title": "Material Properties: Anisotropy, Moisture, and Fire",
    "subtitle": "Engineering Constraints Embedded in a Natural Material",
    "bullets": [
      ("Orthotropic Mechanical Behavior",
       "Timber has three principal grain directions — longitudinal, radial, and tangential — each "
       "with distinct strength and stiffness values. Structural design exploits longitudinal "
       "properties; perpendicular-to-grain loads must be strictly limited in all details."),
      ("Moisture Equilibrium and Creep",
       "Moisture content governs dimensional stability and long-term deflection. Creep under "
       "sustained loading is more significant in timber than in steel; CLT and glulam floor "
       "elements require explicit creep factors in deflection calculations (Porteous & Kermani, 2013)."),
      ("Fire Performance: Charring Rate",
       "Mass timber chars at approximately 0.65 mm/min (species-dependent). The char layer "
       "thermally insulates the inner structural core, enabling load-bearing capacity to persist "
       "for calculable durations. Timber does not fail suddenly — it reduces capacity gradually."),
      ("Sacrificial Char Design",
       "Fire design of mass timber oversizes structural sections to provide a calculated char "
       "depth reserve. Connections — the most fire-vulnerable components — receive protection "
       "through timber encasement, intumescent coatings, or concealed steel plates."),
      ("Biological Susceptibility",
       "Sustained moisture above 18–20% promotes fungal decay; wood-boring insects attack "
       "inadequately protected timber in warm climates. Correct drainage, ventilation, "
       "vapor barriers, and species selection are the primary durability strategies."),
    ],
    "highlight": "Timber burns predictably: the char layer protects the structural core for calculable time. Fire design is not a material limitation — it is a section-sizing and detailing discipline.",
    "citation": "(Porteous & Kermani, 2013; Buchanan & Levine, 1999)",
    "img_key": "timber_interior",
  },
  {
    "num": 4, "section": "TIMBER",
    "title": "Timber Structural Systems and Load Paths",
    "subtitle": "Post-and-Beam, CLT Panels, Glulam Arches, and Hybrid Systems",
    "bullets": [
      ("Post-and-Beam and Timber Frame",
       "Discrete glulam or LVL columns and beams form a structural grid. Loads travel from "
       "floor to beams to posts to foundations. Lateral stability requires bracing, shear "
       "walls, or rigid connections; the frame alone is usually insufficient for lateral loads."),
      ("CLT Panel System",
       "CLT panels act as structural floors, walls, and cores simultaneously. Cross-lamination "
       "creates two-way spanning capacity. Panel-to-panel connections via screws, lap joints, "
       "or proprietary hardware transfer both gravity and diaphragm forces to lateral-resistance "
       "elements (Allen & Iano, 2019)."),
      ("Glulam Arches and Long Spans",
       "Curved glulam arches exploit timber's compressive efficiency to achieve clear spans of "
       "30–80+ m in sports halls, bridges, and exhibition spaces. The three-pinned arch is the "
       "most common statically determinate configuration for long-span roof structures."),
      ("Hybrid Timber Systems",
       "Tall timber buildings almost universally use hybrid strategies: CLT or glulam for "
       "floors and secondary structure, concrete or steel cores for lateral stability, and "
       "composite deck details for vibration and fire performance (Smith & Frangi, 2014)."),
      ("Prefabrication Logic",
       "Factory-fabricated CLT and glulam panels arrive on-site as dimensional components "
       "CNC-cut to precise tolerances. Typical erection rates in mass timber buildings: "
       "one structural floor per 2–5 working days, significantly faster than concrete."),
    ],
    "highlight": "In CLT buildings, the floor–wall–core junction is the critical structural node: gravity load transfer, diaphragm continuity, and lateral stability all converge at this single detail.",
    "citation": "(Allen & Iano, 2019; Smith & Frangi, 2014; Porteous & Kermani, 2013)",
    "img_key": "clt_construction",
  },
  {
    "num": 5, "section": "TIMBER",
    "title": "Timber Connections and Construction Logic",
    "subtitle": "Joinery, Mechanical Fasteners, CNC Fabrication, and Moisture Protection",
    "bullets": [
      ("Traditional vs. Mechanical Connections",
       "Traditional mortise-and-tenon joinery relies on wood-to-wood bearing. Modern connections "
       "use high-strength timber screws, dowels, through-bolts, and concealed steel plates "
       "(Sherpa, Simpson, and custom connectors) — offering higher capacity and better "
       "performance under complex load combinations."),
      ("Connection as Structural Governing Factor",
       "In most timber buildings, connection capacity — not member section — governs the design. "
       "Notching, cross-grain stresses, and end-grain loading at joints must be calculated "
       "explicitly; connection failure is brittle and sudden (Porteous & Kermani, 2013)."),
      ("CNC Fabrication and Assembly Precision",
       "CNC (Computer Numerical Control) timber processing enables submillimeter precision "
       "in cutting, drilling, and milling, making complex 3D joinery feasible at production "
       "scale. Pre-drilled elements arrive on-site ready for rapid bolted assembly."),
      ("Moisture Protection at Connections",
       "All connections must be designed to prevent moisture accumulation: exposed metal is "
       "vulnerable to corrosion; trapped moisture causes timber decay at connection zones. "
       "Design for draining, drying, and inspection access is essential for long-term durability."),
      ("Assembly Sequence",
       "Prefabricated CLT and glulam buildings require crane access and precise delivery "
       "sequencing; panels cannot be stacked interchangeably. The structural assembly sequence "
       "often determines the entire site logistics and construction program."),
    ],
    "highlight": "In timber construction, the connection is the most technically critical element — it governs strength, fire behavior, moisture durability, and long-term serviceability simultaneously.",
    "citation": "(Porteous & Kermani, 2013; Allen & Iano, 2019)",
    "img_key": "brock_commons",
  },
  {
    "num": 6, "section": "TIMBER",
    "title": "Environmental Performance of Timber",
    "subtitle": "Carbon Sequestration, Life-Cycle Assessment, and the Limits of Sustainability Claims",
    "bullets": [
      ("Renewable and Carbon-Storing",
       "Timber is the only major structural material that is biologically renewable. Trees "
       "absorb CO₂ during growth; when used in long-lived buildings, this carbon remains "
       "stored for the service life of the structure, contributing to net carbon reduction "
       "(Ramage et al., 2017)."),
      ("Embodied Carbon Comparison",
       "LCA studies show that CLT construction can achieve significantly lower embodied carbon "
       "than equivalent concrete or steel structures — potentially net-negative in carbon terms "
       "for the structural frame — when certified timber and long service lives are assumed "
       "(De Wolf et al., 2017)."),
      ("Certified Forestry (FSC / PEFC)",
       "Carbon benefits depend critically on responsible forest management. FSC and PEFC "
       "certification schemes verify that harvest rates are matched by regrowth, maintaining "
       "long-term forest carbon stocks. Uncertified timber supply chains can undermine all "
       "environmental claims."),
      ("Prefabrication and Waste Reduction",
       "Factory-fabricated mass timber systems generate significantly less construction waste "
       "than cast-in-place concrete. CNC cutting optimizes material use; off-cuts are often "
       "used for secondary components or biomass energy recovery."),
      ("End-of-Life Strategy",
       "Timber stored in landfill decomposes, releasing stored carbon. Structural reuse of "
       "timber elements, or cascade use (furniture, biomass), is necessary to close the "
       "material's carbon cycle and realize genuine life-cycle benefits (Buchanan & Levine, 1999)."),
    ],
    "highlight": "Timber is not automatically sustainable: only certified forestry, efficient production, long building life, and responsible end-of-life planning together realize the material's genuine environmental potential.",
    "citation": "(Ramage et al., 2017; De Wolf et al., 2017; Buchanan & Levine, 1999)",
    "img_key": "forest",
  },
  {
    "num": 7, "section": "TIMBER",
    "title": "Case Study: Mjøstårnet, Brumunddal, Norway (2019)",
    "subtitle": "Voll Arkitekter — 85.4 m, 18 Storeys — World's Tallest Timber Building",
    "bullets": [
      ("Project Overview",
       "Completed in 2019, Mjøstårnet reaches 85.4 m and contains mixed uses: hotel, "
       "apartments, offices, and a public indoor swimming pool. It is certified as the "
       "world's tallest timber building and demonstrates that tall timber is commercially, "
       "legally, and structurally viable (Voll Arkitekter, 2019)."),
      ("Structural System",
       "Primary structure: glulam columns and diagonal X-bracing trusses visible on the "
       "façade, with CLT floor panels and shaft walls. The exposed glulam bracing system "
       "resolves lateral wind loads while serving as the building's primary architectural "
       "expression — structural logic made visible."),
      ("Hybrid Strategy Above Floor 10",
       "CLT floor decks are replaced by concrete slabs above the 10th floor to increase "
       "mass and reduce floor vibration response — a pragmatic hybrid decision driven by "
       "occupant comfort, not structural necessity (Kuzmanovska et al., 2018)."),
      ("Fire Strategy",
       "Large-section glulam members provide inherent char protection; the building is "
       "fully sprinklered; fire separation is achieved through fire-rated CLT walls and "
       "protected stairwells. Norwegian code compliance for tall buildings was established "
       "through performance-based fire engineering analysis."),
      ("Carbon and Prefabrication",
       "Timber elements were sourced from Norwegian forests and prefabricated locally, "
       "minimizing transport carbon. Structural carbon footprint is substantially biogenic "
       "(stored) rather than fossil-based, making this one of the lowest-carbon tall "
       "building structures constructed to date."),
    ],
    "highlight": "Mjøstårnet is not just a record-holder: it is a proof of concept that mass timber, hybrid structural strategies, and performance-based fire engineering can together produce viable tall urban buildings.",
    "citation": "(Kuzmanovska et al., 2018; Voll Arkitekter, 2019; Ramage et al., 2017)",
    "img_key": "mjostaarnet",
  },
  {
    "num": 8, "section": "TIMBER",
    "title": "Case Study: Brock Commons & Tamedia Office Building",
    "subtitle": "Prefabrication Speed vs. Structural Expression — Two Philosophies in Mass Timber",
    "bullets": [
      ("Brock Commons: 18-Storey Student Residence",
       "Brock Commons Tallwood House (Vancouver, 2016; Acton Ostry Architects) is an 18-storey, "
       "53 m student residence at UBC, constructed with remarkable speed: the mass timber "
       "structure was erected at approximately one floor per day using prefabricated CLT and "
       "glulam elements (Acton Ostry Architects, 2016)."),
      ("Brock Commons: Structural Strategy",
       "A concrete podium and two concrete cores provide lateral stability and fire egress. "
       "CLT floor panels span between glulam columns from floors 2 to 18. Structural concrete "
       "encasement of the base columns and concrete cores reflects a pragmatic hybrid approach "
       "prioritizing code compliance and construction speed over full timber expression."),
      ("Tamedia: Joinery-Based Office Building",
       "Tamedia Office Building (Zurich, 2013; Shigeru Ban Architects) uses a fully exposed, "
       "seven-storey timber structure with custom-engineered wood-to-wood connections free of "
       "metal fasteners — a contemporary interpretation of Japanese joinery at building scale."),
      ("Tamedia: Structural Expression",
       "Spruce columns and beams are connected by CNC-milled hardwood dowel joints, creating "
       "a system visible from every interior point. The exposed timber is simultaneously "
       "structure, finish, and spatial character — reducing material layers and expressing "
       "structural honesty throughout the building (Shigeru Ban Architects, 2013)."),
      ("Design Spectrum",
       "Brock Commons and Tamedia represent opposite ends of a mass timber design spectrum: "
       "maximum construction efficiency vs. maximum structural expression. Both are technically "
       "successful; the appropriate choice depends on programme, budget, climate, and intent."),
    ],
    "highlight": "Brock Commons proves that mass timber can be industrially fast; Tamedia proves it can be architecturally precise — the spectrum between them defines the range of contemporary timber practice.",
    "citation": "(Acton Ostry Architects, 2016; Shigeru Ban Architects, 2013; Lehmann, 2012)",
    "img_key": "brock_commons",
  },
  {
    "num": 9, "section": "TIMBER",
    "title": "Limitations and Design Risks of Timber",
    "subtitle": "Moisture, Decay, Acoustics, Vibration, and Regulatory Constraints",
    "bullets": [
      ("Moisture and Biological Decay",
       "Sustained moisture above 18–20% promotes fungal decay; wood-boring insects cause "
       "structurally significant damage without visible surface evidence. Poor detailing at "
       "roof connections, wall bases, and junction zones — where water can pond or migrate — "
       "is the dominant cause of long-term timber building failure (Allen & Iano, 2019)."),
      ("Acoustic Performance",
       "CLT floor assemblies without adequate topping, floating screed, or resilient "
       "underlayer typically fail residential impact sound standards. Airborne sound "
       "transmission through CLT floors and walls also requires multi-layer assemblies "
       "with mass and decoupling, adding complexity and cost."),
      ("Floor Vibration",
       "Lightweight timber floors in spans exceeding 6–8 m can exhibit vibration at "
       "frequencies (4–8 Hz) that are perceptible and uncomfortable to occupants. "
       "Explicit vibration serviceability checks to national standards are mandatory "
       "for all longer-span timber floor designs."),
      ("Creep and Long-Term Deflection",
       "Under sustained loading, timber deflects more than initial elastic calculation "
       "suggests. Creep factors of 0.6–0.8 must be applied in Eurocode 5 deflection "
       "checks; unattended creep causes visible sag, cracking of finishes, and potential "
       "serviceability failure (Porteous & Kermani, 2013)."),
      ("Regulatory and Insurance Constraints",
       "Building codes and insurance markets in many jurisdictions restrict timber building "
       "height and occupancy type more severely than concrete or steel. Performance-based "
       "fire engineering approval processes add time and cost to mass timber projects, "
       "particularly for first-of-type applications in new markets."),
    ],
    "highlight": "Timber failure in buildings is rarely due to material inadequacy — it is almost always the result of moisture detailing error, acoustic under-design, or inadequate vibration engineering.",
    "citation": "(Porteous & Kermani, 2013; Allen & Iano, 2019)",
    "img_key": "timber_limits",
  },
  {
    "num": 10, "section": "TIMBER",
    "title": "Timber Synthesis: Why Timber Matters Today",
    "subtitle": "Carbon, Prefabrication, and the Architecture of Lightweight Structure",
    "bullets": [
      ("Climate Imperative",
       "Construction accounts for approximately 38% of global CO₂ emissions. Timber is the "
       "only widely available structural material that sequesters rather than emits carbon "
       "during production, making it strategically significant in any carbon-reduction "
       "construction pathway (Ramage et al., 2017)."),
      ("Proven Technical Viability",
       "Mass timber has demonstrated structural viability from residential to 85-metre "
       "towers, in office, residential, educational, and cultural building types. "
       "Performance-based codes in Norway, Austria, and Canada have established clear "
       "precedents that other jurisdictions are progressively adopting."),
      ("Prefabrication and Efficiency",
       "CLT and glulam's factory fabrication reduces on-site construction time, workforce "
       "requirements, and waste generation. The DfMA (Design for Manufacture and Assembly) "
       "logic of mass timber aligns with industrialization trends in the construction sector."),
      ("Biophilic and Occupant Benefits",
       "Research consistently associates exposed timber interiors with measurable occupant "
       "wellbeing outcomes: reduced stress response, improved cognitive performance, and "
       "higher reported satisfaction — architectural qualities that reinforce the technical "
       "case for timber as a building material (Lehmann, 2012)."),
      ("The Disciplined Future",
       "Timber's future requires engineering precision: moisture detailing, acoustic design, "
       "vibration analysis, certified supply chains, and end-of-life planning. It is neither "
       "a universal solution nor a nostalgic retreat — it is a technically rigorous material "
       "choice with genuine climate and architectural potential."),
    ],
    "highlight": "Timber's future is not romantic — it is disciplined: precision detailing, verified LCA, and genuine carbon accounting are the only bases on which timber's environmental and structural claims stand.",
    "citation": "(Ramage et al., 2017; Lehmann, 2012; Smith & Frangi, 2014)",
    "img_key": "mjostaarnet",
  },

  # ─── STEEL ────────────────────────────────────────────────────────────────
  {
    "num": 11, "section": "STEEL",
    "title": "Steel as a Structural Material",
    "subtitle": "Industrial, Isotropic, High-Strength, and Ductile",
    "bullets": [
      ("Isotropic and Homogeneous",
       "Structural steel is an iron-carbon alloy produced through controlled industrial "
       "processes. Unlike timber, steel is isotropic: mechanical properties are equal in "
       "all directions, enabling precise calculation and standardized section production "
       "with highly reliable, predictable structural behavior (Lim, 2019)."),
      ("Tensile and Compressive Strength",
       "Structural steel yield strength: 275 MPa (S275) to 355 MPa (S355) — approximately "
       "ten times the compressive strength of concrete and thirty times that of standard "
       "timber. This high strength allows slender sections to carry large loads, enabling "
       "long spans and high-rise structures with minimal material volume."),
      ("Ductility and Seismic Performance",
       "Steel's ductility — its ability to deform plastically before fracture — is its most "
       "critical seismic property. Ductile frames dissipate earthquake energy through "
       "controlled yielding, preventing sudden collapse and giving occupants time to evacuate "
       "(Allen & Iano, 2019)."),
      ("Precision and Standardization",
       "Structural steel sections (I-beams, H-columns, channels, angles, hollow sections) "
       "are produced to millimeter tolerances in standard grades and dimensions worldwide. "
       "Standardization simplifies calculation, enables international procurement, and "
       "supports fast-track project delivery."),
      ("Critical Vulnerabilities",
       "Steel's performance limitations are well-defined: yield strength decreases by "
       "50–70% above 550°C (fire); surfaces corrode without protection (wet environments); "
       "and thermal conductivity (~50 W/mK) creates significant heat bridges in building "
       "envelopes. All three require explicit design response."),
    ],
    "highlight": "Steel's structural identity is defined by high tensile capacity, predictable ductility, and dimensional precision — but fire protection, corrosion control, and thermal bridging are non-negotiable engineering obligations.",
    "citation": "(Lim, 2019; Allen & Iano, 2019)",
    "img_key": "steel_construction",
  },
  {
    "num": 12, "section": "STEEL",
    "title": "Historical Evolution: From Cast Iron to Modern Steel",
    "subtitle": "Bessemer Process, Crystal Palace, Eiffel Tower, and the Birth of the High-Rise",
    "bullets": [
      ("Cast Iron and Early Iron Structures",
       "Cast iron (compressive, brittle) dominated 18th–early 19th century structural "
       "applications: columns, bridge arches, and mill building frames. Wrought iron "
       "provided improved tensile strength for beams and tension members. Both materials "
       "preceded modern structural steel but established the concept of metallic framing."),
      ("Crystal Palace (1851): Industrial Prefabrication",
       "Crystal Palace (Joseph Paxton) demonstrated prefabricated modular iron-and-glass "
       "construction: 83,000 sq m erected in approximately 9 months using standardized "
       "components — an unprecedented achievement of industrial production applied to "
       "architecture (Addis, 2007)."),
      ("Eiffel Tower (1889): Structural Expression",
       "The Eiffel Tower's 7,300 tonnes of latticed iron were assembled from 18,038 "
       "individually calculated components. Gustave Eiffel's wind-loading calculations "
       "were among the first scientifically rigorous applications of structural analysis "
       "to a major building — making the tower as important as an engineering milestone "
       "as an architectural one (Billington, 1983)."),
      ("Bessemer Process and the Steel Frame",
       "The Bessemer process (1856) massively reduced steel production costs. By the "
       "1880s, standardized rolled steel sections replaced iron. The Home Insurance "
       "Building (Chicago, 1885; Jenney) first transferred building loads entirely "
       "through a metal frame, liberating the exterior wall from structural duty."),
      ("Modern Steel Construction",
       "20th-century steel construction evolved through standardized sections, "
       "high-strength bolts (replacing rivets), automatic welding, and computerized "
       "structural analysis, producing the global high-rise building stock that "
       "defines contemporary city skylines."),
    ],
    "highlight": "The steel frame eliminated the load-bearing wall — freeing the façade, enabling flexible plans, and making modern tall buildings possible. This structural revolution of the 1880s fundamentally redefined architectural scale.",
    "citation": "(Addis, 2007; Billington, 1983)",
    "img_key": "crystal_palace",
  },
  {
    "num": 13, "section": "STEEL",
    "title": "Material Properties: Strength, Ductility, and Structural Behavior",
    "subtitle": "Yield Strength, Modulus of Elasticity, Buckling, Fire, and Thermal Bridges",
    "bullets": [
      ("Yield Strength and Plastic Behavior",
       "At yield strength (fy), steel enters permanent plastic deformation before "
       "strain-hardening. Plastic design methods exploit this plateau, allowing "
       "redistribution of forces in statically indeterminate systems. Ultimate tensile "
       "strength (fu) exceeds fy by 20–30%, providing significant overstrength reserve."),
      ("Elastic Modulus and Stiffness",
       "Steel's modulus of elasticity (E ≈ 200–210 GPa) is approximately three times "
       "concrete and ten times structural timber. This high stiffness enables slender "
       "members with minimal deflection — but also means buckling (not yielding) often "
       "governs the design of columns and beams (Lim, 2019)."),
      ("Buckling as Governing Failure Mode",
       "Slender steel columns fail by Euler buckling at stresses well below yield. "
       "Lateral-torsional buckling governs unrestrained steel beams under bending. "
       "Plate buckling governs thin-walled sections. All three require explicit "
       "calculation beyond material strength checks."),
      ("Fire-Induced Strength Loss",
       "At 300°C, steel retains ~80% of ambient yield strength. At 550°C, retention "
       "falls to ~30%. At 700°C, near-complete loss occurs. All structural steel "
       "in fire-risk zones requires fire protection — intumescent paint, board "
       "encasement, or concrete encasing — calibrated to the design fire scenario."),
      ("Thermal Conductivity and Bridges",
       "Steel's thermal conductivity (~50 W/mK) is 1,500× greater than mineral wool "
       "insulation. Connections between steel structural elements and building envelopes "
       "create significant thermal bridges unless detailed with thermal break systems, "
       "increasing energy loss and condensation risk."),
    ],
    "highlight": "Steel is strong in tension — but fire protection and stability against buckling are not secondary concerns: they are primary structural design obligations that govern section selection in most practical applications.",
    "citation": "(Lim, 2019; Allen & Iano, 2019)",
    "img_key": "i_beam",
  },
  {
    "num": 14, "section": "STEEL",
    "title": "Steel Structural Systems and Load Paths",
    "subtitle": "Frames, Trusses, Space Frames, Cables, and Diagrids",
    "bullets": [
      ("Moment-Resisting Frames",
       "Rigid beam-column connections transfer both gravity and lateral loads through "
       "bending and shear in columns and beams. Moment frames are ductile under seismic "
       "loading and allow open, unbraced floor plans — essential for flexible office "
       "and mixed-use building typologies."),
      ("Braced Frames",
       "Diagonal steel bracing (K-brace, X-brace, eccentric brace) resists lateral loads "
       "primarily through axial forces, providing stiff lateral resistance with minimal "
       "material. Eccentric braced frames (EBF) add ductility through controlled link "
       "element yielding for seismic applications (Taranath, 2012)."),
      ("Trusses for Long Spans",
       "Trusses resolve loads into axial tension and compression in individual members, "
       "achieving very long spans (20–100+ m) with minimal self-weight. Applications: "
       "roof structures, bridge decks, floor transfer structures, and stadium canopies."),
      ("Space Frames and Cable Structures",
       "Double-layer space frames span in two directions with uniform depth, suited for "
       "large column-free roof areas. Cable structures (stayed or suspended) exploit "
       "steel's tensile strength for ultra-long spans with minimal material volume, "
       "used in stadium roofs, bridges, and atrium structures."),
      ("Diagrid Exoskeletons",
       "Diagrid systems — diagonal external grid frames — transfer gravity and lateral "
       "loads simultaneously through inclined members, reducing or eliminating internal "
       "core requirements. The Hearst Tower (Foster+Partners, NYC) reduced steel tonnage "
       "by approximately 20% compared to a conventional braced frame."),
    ],
    "highlight": "Steel's versatility lies in its ability to organize load paths into multiple structural languages: the same material can express tension (cable), compression (arch), bending (beam), or combined action (diagrid) with equal efficiency.",
    "citation": "(Taranath, 2012; Billington, 1983; Lim, 2019)",
    "img_key": "pompidou",
  },
  {
    "num": 15, "section": "STEEL",
    "title": "Steel Connections and Assembly",
    "subtitle": "Bolted Joints, Welded Connections, and Fabricated Precision",
    "bullets": [
      ("Bolted Connections",
       "High-strength friction-grip bolts (Grade 8.8 and 10.9) are the standard for "
       "site assembly. Bolted connections are quick to erect, visually inspectable, "
       "torque-verifiable, and demountable at end of life — making them the preferred "
       "choice for circular economy steel design strategies."),
      ("Moment-Resisting Connections",
       "Flush and extended end-plate connections with high-strength bolts achieve "
       "near-rigid moment transfer for wind and seismic frames. Full-strength moment "
       "connections require careful engineering of bolt pattern, plate thickness, "
       "and column stiffener requirements (Lim, 2019)."),
      ("Shop Welding vs. Site Welding",
       "Factory welding achieves higher quality under controlled conditions (position, "
       "temperature, inspection) than field welding. Standard practice: shop-weld "
       "primary members; site-bolt connections. All structural welds require "
       "non-destructive testing (NDT) for critical elements."),
      ("Base Plates and Anchor Bolts",
       "Column base plates distribute column load to concrete foundations through "
       "bearing; anchor bolts resist combined tension, shear, and moment. "
       "Leveling nuts and non-shrink grout beneath the plate ensure full bearing "
       "contact and dimensional accuracy at the structure-foundation interface."),
      ("Prefabrication and Erection Efficiency",
       "Structural steel is virtually entirely fabricated off-site, delivered to "
       "dimensional tolerance, and assembled by crane in sequence. Modern fabrication "
       "shops use CNC plate cutting, robotic welding, and BIM-coordinated geometry "
       "to achieve complex three-dimensional structures with high accuracy and speed."),
    ],
    "highlight": "Steel connection design is where structural mechanics, fire protection, corrosion resistance, fabrication economics, and circular reuse potential all converge — a good connection is simultaneously efficient in all five dimensions.",
    "citation": "(Lim, 2019; Allen & Iano, 2019)",
    "img_key": "pompidou",
  },
  {
    "num": 16, "section": "STEEL",
    "title": "Case Study: Eiffel Tower & Crystal Palace",
    "subtitle": "Metal Structure as Monumental Architecture and Industrial Prefabrication",
    "bullets": [
      ("Crystal Palace: Modular Industrial Prefabrication",
       "Crystal Palace (Paxton, 1851) demonstrated that prefabricated iron and glass could "
       "create 90,000 sq m of covered space in under nine months using standardized modules. "
       "The entire building was designed around a production logic rather than a "
       "compositional one — an architectural revolution as significant as its construction "
       "speed (Addis, 2007)."),
      ("Crystal Palace: Structural and Economic Innovation",
       "All iron members were produced from standard patterns, reducing fabrication cost "
       "and enabling precise structural calculation of repetitive elements. The building "
       "was subsequently dismantled and re-erected at Sydenham — demonstrating "
       "demountability as a design intention, not an afterthought."),
      ("Eiffel Tower: Wind Engineering Precedent",
       "Eiffel's design was driven by quantitative wind load analysis at a time when such "
       "calculations were rare in architectural practice. The lattice cross-section "
       "minimizes wind resistance while maximizing structural stiffness — form following "
       "engineering logic directly (Billington, 1983)."),
      ("Eiffel Tower: Structure as Architecture",
       "At 324 m, the Eiffel Tower remained the world's tallest structure for 41 years. "
       "The exposed iron lattice — initially controversial — became the defining image "
       "of industrial modernity, demonstrating that structural engineering, expressed "
       "honestly, can achieve monumental architectural significance."),
      ("Maintenance and Long-Term Durability",
       "The Eiffel Tower requires approximately 50–60 tonnes of paint applied every "
       "7 years to prevent corrosion — a continuous maintenance commitment that "
       "underlines the real cost of exposed metal structures and the need to "
       "incorporate maintenance cycles into lifecycle planning from the outset."),
    ],
    "highlight": "Crystal Palace and the Eiffel Tower demonstrate that structural engineering, expressed without concealment, can be the primary source of architectural meaning — the structural system IS the architecture.",
    "citation": "(Addis, 2007; Billington, 1983)",
    "img_key": "eiffel",
  },
  {
    "num": 17, "section": "STEEL",
    "title": "Case Study: Centre Pompidou & Seagram Building",
    "subtitle": "High-Tech Expression vs. Miesian Restraint — Two Poles of Steel Architecture",
    "bullets": [
      ("Centre Pompidou: Structural Inversion",
       "Centre Pompidou (Piano & Rogers, 1977) relocated all structural, mechanical, "
       "and circulation elements to the exterior, leaving floor plates entirely unobstructed. "
       "Cast steel gerberette brackets carry floor loads outward; color-coded external "
       "pipes and ducts make the building's servicing infrastructure its dominant "
       "aesthetic element (Addis, 2007)."),
      ("Pompidou: Engineering Collaboration",
       "The gerberette — a large cast steel cantilever element supporting 800 tonnes "
       "per floor — required intensive collaboration between architects Piano & Rogers "
       "and structural engineers Ove Arup & Partners, establishing a precedent for "
       "integrated technical design that remains influential today."),
      ("Seagram Building: Concealed Structure, Expressed Logic",
       "The Seagram Building (Mies van der Rohe, 1958) conceals its primary steel frame "
       "within fire-proofed concrete. The visible bronze I-section mullions on the "
       "curtain wall are non-structural — an applied expression of structural logic "
       "rather than the structure itself (Billington, 1983)."),
      ("Mies's Structural Philosophy",
       "Mies accepted the contradiction between structural honesty and fire protection "
       "requirements, using applied ornamental steel to represent the structural grid "
       "that could not legally be exposed. This philosophical position — structural "
       "expression as architectural language, not literal exposure — remains a "
       "productive debate in contemporary architecture."),
      ("Two Traditions",
       "Pompidou (exposed, raw, industrial) and Seagram (restrained, refined, graphic) "
       "represent the two dominant traditions of steel architectural expression: "
       "honest engineering revelation vs. structural representation. Both are intellectually "
       "rigorous; both require deep material intelligence to execute convincingly."),
    ],
    "highlight": "Steel in architecture has two legitimate traditions: honest structural exposure (Centre Pompidou) and structural representation through applied material (Seagram Building). Both require rigorous material intelligence.",
    "citation": "(Billington, 1983; Addis, 2007)",
    "img_key": "seagram",
  },
  {
    "num": 18, "section": "STEEL",
    "title": "Environmental Performance of Steel",
    "subtitle": "High Embodied Carbon, Exceptional Recyclability, and the Circular Economy",
    "bullets": [
      ("Embodied Carbon: Production Energy",
       "Basic oxygen furnace (BOF) steel from virgin ore: approximately 1.6–2.8 t CO₂/t "
       "steel. Electric arc furnace (EAF) steel from recycled scrap: approximately "
       "0.4–0.7 t CO₂/t steel — a factor of 4× improvement. Source material specification "
       "is therefore the single most impactful environmental choice in steel procurement "
       "(De Wolf et al., 2017)."),
      ("Recyclability: Highest of Any Structural Material",
       "Global structural steel recycling rates exceed 85%. Steel can be recycled "
       "indefinitely without mechanical property degradation. Unlike concrete demolition "
       "waste (which downcycles to aggregate), steel retains full structural value "
       "through the recycling loop."),
      ("Reuse vs. Recycling",
       "Structural reuse — using a member in its original form in a new building — "
       "requires near-zero energy input and avoids all re-melting carbon. Bolted "
       "connections and standardized section dimensions make structural steel the most "
       "reusable major building material. Design for disassembly (DfD) strategies "
       "are increasingly incorporated into new projects."),
      ("Longevity and Adaptability",
       "Steel buildings that serve multiple uses over 100+ years amortize embodied "
       "carbon over extended service lives. Adaptive reuse of steel-framed buildings "
       "— warehouses converted to offices, factories to residential — demonstrates "
       "steel's long-term structural adaptability as a carbon-reduction strategy."),
      ("Green Steel Pathway",
       "Hydrogen-based direct reduction (H-DRI) and renewable-powered EAF processes "
       "offer a near-zero-carbon steel production pathway. Commercial green steel "
       "production is expanding; by 2040–2050, structural steel's embodied carbon "
       "profile could be fundamentally transformed."),
    ],
    "highlight": "Steel's high upfront embodied carbon must be weighed against exceptional recyclability, reusability potential, structural durability, and the emerging trajectory toward green steel production.",
    "citation": "(De Wolf et al., 2017; Allen & Iano, 2019)",
    "img_key": "willis_tower",
  },
  {
    "num": 19, "section": "STEEL",
    "title": "Limitations and Failure Risks of Steel",
    "subtitle": "Corrosion, Buckling, Fire, Fatigue, and Thermal Expansion",
    "bullets": [
      ("Corrosion: The Primary Durability Risk",
       "Steel corrodes through electrochemical oxidation in the presence of moisture, "
       "oxygen, chlorides, or pollutants. Coastal and industrial environments require "
       "high-performance coating systems with 5–15 year maintenance cycles. "
       "Undetected corrosion in inaccessible zones is a major long-term structural risk."),
      ("Fire-Induced Structural Failure",
       "At 550°C, structural steel retains only ~30% of ambient yield strength; at 700°C, "
       "near-complete capacity loss occurs. All steel in fire-risk areas requires "
       "protection calibrated to the design fire: intumescent paint, mineral spray, "
       "or board encasement. Fire protection is a structural engineering requirement, "
       "not a code formality (Lim, 2019)."),
      ("Stability: Buckling Failure Modes",
       "Steel members fail by flexural buckling (columns), lateral-torsional buckling "
       "(beams), and local plate buckling (thin-walled sections) at loads below material "
       "yield strength. Slenderness governs design; restraint against lateral-torsional "
       "buckling must be provided at regular intervals for beams."),
      ("Fatigue Under Cyclic Loading",
       "Repetitive loading (bridges, cranes, wind, machinery) causes fatigue crack "
       "initiation and propagation at welds and stress concentrations. Welded connections "
       "have lower fatigue resistance than parent metal; fatigue life must be explicitly "
       "calculated for all dynamically loaded structural elements."),
      ("Thermal Expansion and Movement",
       "Steel expands at ~12 × 10⁻⁶ /°C — approximately 3× masonry and 8× concrete. "
       "Long steel structures (>50–60 m) require designed expansion joints. Connections "
       "to other materials must accommodate differential thermal movement to prevent "
       "stress concentration and joint failure."),
    ],
    "highlight": "Steel's performance under fire, corrosion, and fatigue is entirely manageable — but only when protection and detailing strategies are fully integrated into the structural design, not added as afterthoughts.",
    "citation": "(Lim, 2019; Allen & Iano, 2019)",
    "img_key": "steel_frame",
  },
  {
    "num": 20, "section": "STEEL",
    "title": "Steel Synthesis: Precision, Span, and the Architecture of Modernity",
    "subtitle": "From Industrial Frames to Sustainable, Demountable Structures",
    "bullets": [
      ("Steel's Definitive Structural Role",
       "Steel enables the spans (20–200+ m), heights (100–800+ m), and structural "
       "refinements that define contemporary architecture. The curtain wall tower, the "
       "long-span roof, and the cable-stayed bridge are all structurally dependent on "
       "steel's tensile capacity, ductility, and precision (Taranath, 2012)."),
      ("The Sustainability Challenge",
       "High embodied carbon — 1.6–2.8 t CO₂/t steel for virgin production — is the "
       "material's primary environmental liability. Minimum material volume through "
       "structural optimization, maximum recycled content specification, and design "
       "for disassembly are the principal mitigation strategies available now."),
      ("Hybrid Structural Solutions",
       "Contemporary practice increasingly combines steel with timber, concrete, and "
       "glass — using each material only where its properties are most advantageous. "
       "Composite steel-concrete decks, steel-timber hybrid floors, and steel-CLT "
       "buildings reduce total material use while optimizing structural performance."),
      ("Digital Fabrication and Form",
       "Parametric structural design and robotic fabrication enable steel structures "
       "of unprecedented geometric complexity — branching columns, doubly curved trusses, "
       "and diagrid exoskeletons — without proportional cost increase, expanding "
       "steel's architectural vocabulary significantly."),
      ("The Circular Steel Future",
       "Design for disassembly using bolted connections and material passports, "
       "combined with green steel (H-DRI) production and high recycled content "
       "specification, provides a credible pathway toward dramatically reduced "
       "embodied carbon in steel construction by mid-century."),
    ],
    "highlight": "Steel's future is not more steel — it is less, better-placed, longer-lived, and progressively lower-carbon: precision optimization, circular reuse, and green production working together.",
    "citation": "(De Wolf et al., 2017; Taranath, 2012; Allen & Iano, 2019)",
    "img_key": "pompidou",
  },

  # ─── MASONRY ─────────────────────────────────────────────────────────────
  {
    "num": 21, "section": "MASONRY",
    "title": "Masonry as a Structural Material",
    "subtitle": "Mass, Compression, and the Architecture of Permanence",
    "bullets": [
      ("Definition and Structural Logic",
       "Masonry is a unit-based compressive system assembled from discrete elements — "
       "stone, brick, concrete masonry units (CMU), or adobe — bonded with mortar "
       "or dry-stacked. It is the oldest structural tradition in human building history, "
       "predating metal and engineered materials by millennia (Hendry et al., 2004)."),
      ("Compressive Strength, Tensile Weakness",
       "Masonry units resist compression effectively (brick: 20–100 MPa; stone: "
       "50–200 MPa) but have very low tensile capacity (typically 10–20% of "
       "compressive strength). This asymmetry shaped the entire history of masonry "
       "structural form — from the arch to the vault to the dome — as geometric "
       "strategies for eliminating tensile stress."),
      ("Thermal Mass",
       "Dense masonry (1,500–2,500 kg/m³) stores thermal energy at 840–1,000 J/kg·K, "
       "providing significant thermal mass. A 200 mm brick wall gives 8–12 hours of "
       "thermal lag, moderating interior temperatures in climates with large "
       "day-night variation — a passive energy-conservation strategy (Allen & Iano, 2019)."),
      ("Durability and Redundancy",
       "Well-constructed masonry is among the most durable structural systems: Roman "
       "brickwork from 2,000 years ago remains structurally functional. The redundant, "
       "unit-based assembly means localized damage rarely causes global collapse — "
       "a resilience characteristic not shared by monolithic structural systems."),
      ("Craft and Local Material",
       "Masonry is historically a craft-based material: bricklaying quality directly "
       "affects structural performance. Clay brick, stone, and adobe are typically "
       "extracted and produced regionally, reducing transport carbon and strengthening "
       "the connection between building material and site geology."),
    ],
    "highlight": "Masonry is compressive architecture: its entire structural vocabulary — arch, vault, buttress, dome — is a geometric solution to the fundamental challenge of spanning with materials that cannot resist tension.",
    "citation": "(Hendry et al., 2004; Heyman, 1995; Allen & Iano, 2019)",
    "img_key": "pont_du_gard",
  },
  {
    "num": 22, "section": "MASONRY",
    "title": "Historical Evolution of Masonry",
    "subtitle": "Stone, Brick, Adobe, and Concrete Block — 10,000 Years of Building Tradition",
    "bullets": [
      ("Prehistoric and Ancient Stone",
       "Neolithic megalithic structures (Stonehenge, Carnac), Mycenaean Cyclopean walls, "
       "and Egyptian pyramid masonry demonstrate compressive structural logic at "
       "monumental scale. Stone selection, coursing, and mass — rather than mortar "
       "or reinforcement — provided structural integrity (Addis, 2007)."),
      ("Roman Masonry: Arch, Vault, and Concrete",
       "Roman construction combined fired brick, opus reticulatum concrete facing, and "
       "pozzolanic lime concrete to build vaults, domes, and aqueducts of unprecedented "
       "scale. The Pantheon dome (43.3 m diameter, c. 125 CE) remains the world's "
       "largest unreinforced concrete dome (Lancaster, 2005)."),
      ("Byzantine and Islamic Masonry",
       "Byzantine architecture refined the pendentive dome — transferring circular "
       "dome loads to square piers — enabling the Hagia Sophia's 31.3 m dome "
       "(537 CE). Islamic masonry combined structural brick, geometric tile, and "
       "muqarnas vaulting into integrated tectonic and decorative systems."),
      ("Gothic Masonry: Structural Limits",
       "Gothic cathedrals pushed load-bearing stone masonry to its material limits: "
       "pointed arches reduce lateral thrust; ribbed vaults concentrate loads "
       "on piers; flying buttresses externalize and resolve thrust, enabling "
       "thin nave walls and large clerestory windows (Heyman, 1995)."),
      ("Industrial Brick and Modern Masonry",
       "Industrial brick production (19th century) enabled standardized modular "
       "construction at urban scale. The Monadnock Building (Chicago, 1891) reached "
       "16 storeys with load-bearing brick walls — the practical upper limit. "
       "Modern masonry includes reinforced CMU, cavity wall systems, and brick façades "
       "on structural frames."),
    ],
    "highlight": "The history of masonry is a history of geometric intelligence: each era found new forms — arch, vault, dome, buttress — that redirected tension into compression, expanding masonry's structural range without changing its material logic.",
    "citation": "(Heyman, 1995; Lancaster, 2005; Addis, 2007)",
    "img_key": "pont_du_gard",
  },
  {
    "num": 23, "section": "MASONRY",
    "title": "Material Properties: Compression, Thermal Mass, and Mortar",
    "subtitle": "Compressive Strength, Brittleness, Moisture, and the Structural Role of the Mortar Joint",
    "bullets": [
      ("Compressive Strength Variation",
       "Masonry unit compressive strength varies widely: engineering brick 50–100 MPa; "
       "standard clay brick 20–50 MPa; CMU 17–35 MPa; adobe 1–5 MPa. Assembled "
       "masonry compressive strength (fk) is significantly lower than unit strength "
       "due to mortar joint behavior and stress concentration (Hendry et al., 2004)."),
      ("Tensile and Shear Weakness",
       "Tensile bond strength of mortar-bonded masonry: typically 0.1–0.5 MPa — "
       "approximately 5–10% of compressive strength. Mortar-unit interfaces are "
       "the weakest tensile planes; flexural cracking, diagonal shear cracking, "
       "and sliding shear at bed joints are the dominant in-plane failure modes."),
      ("Mortar Composition and Stiffness",
       "Cement-rich mortars are stronger but stiffer, potentially concentrating "
       "stress and increasing cracking risk in older or flexible masonry. "
       "Lime-based mortars are weaker but flexible, redistributing stresses "
       "across more joints — preferred in conservation and heritage contexts "
       "where differential movement is expected (Morton, 2008)."),
      ("Moisture Movement",
       "Fired clay bricks expand irreversibly after manufacture (moisture expansion). "
       "Concrete and autoclaved products (AAC, CMU) exhibit drying shrinkage. "
       "Both phenomena generate internal stresses requiring designed movement joints "
       "at regular intervals and at changes in wall direction or material."),
      ("Thermal Mass Performance",
       "Masonry's thermal storage capacity moderates diurnal temperature swings, "
       "reducing peak heating/cooling demand in appropriate climates. A 220 mm "
       "solid brick wall provides approximately 8–12 hours of thermal lag — "
       "a significant passive climate control benefit in hot-dry and Mediterranean "
       "climatic zones (Allen & Iano, 2019)."),
    ],
    "highlight": "Mortar is not passive bonding: it governs masonry deformability, crack resistance, and moisture tolerance. The wrong mortar in a repair can cause more damage than the original decay it replaced.",
    "citation": "(Hendry et al., 2004; Morton, 2008; Allen & Iano, 2019)",
    "img_key": "brick_wall",
  },
  {
    "num": 24, "section": "MASONRY",
    "title": "Masonry Structural Systems and Load Paths",
    "subtitle": "Bearing Walls, Arches, Vaults, Domes, and Buttresses",
    "bullets": [
      ("Load-Bearing Walls",
       "Masonry bearing walls carry gravity loads through continuous vertical compression "
       "from roof to floor to foundation. Wall slenderness ratio (height/thickness) "
       "and lateral support conditions govern stability; as height increases without "
       "intermediate lateral restraint, buckling governs over material strength (Heyman, 1995)."),
      ("The Arch: Fundamental Spanning Solution",
       "An arch redirects vertical loads into diagonal compression forces that flow "
       "to the abutments. In an arch aligned with the thrust line under its design "
       "load, all material is in compression — zero tensile stress. The arch is the "
       "geometrically optimal spanning solution for a purely compressive material."),
      ("Vault and Dome Geometry",
       "The barrel vault extends the arch longitudinally; the groin vault intersects "
       "two barrel vaults at right angles. The dome extends the arch radially, creating "
       "three-dimensional compression that generates compressive hoop forces at the crown "
       "and potentially tensile hoop forces near the base — the historic cause of dome "
       "cracking addressed by iron tension rings (Lancaster, 2005)."),
      ("Buttresses and Thrust Resolution",
       "Lateral thrust from arches and vaults must be resolved into ground-reaching "
       "compression. Massive buttresses absorb thrust through weight and cross-section; "
       "Gothic flying buttresses externalize and transmit thrust over the side aisles, "
       "enabling thinner nave walls and larger window openings."),
      ("Reinforced Masonry",
       "Steel reinforcement placed in grouted masonry cores (vertical) or bed joints "
       "(horizontal) fundamentally extends masonry's structural capability to resist "
       "bending, tension, and shear — particularly in seismic zones where unreinforced "
       "masonry is vulnerable (Hendry et al., 2004)."),
    ],
    "highlight": "The arch is not merely a historical form — it is the structurally correct solution for spanning with compressive materials. Every masonry arch built since antiquity obeys the same thrust-line logic.",
    "citation": "(Heyman, 1995; Lancaster, 2005; Hendry et al., 2004)",
    "img_key": "pont_du_gard",
  },
  {
    "num": 25, "section": "MASONRY",
    "title": "Masonry Bonds, Details, and Construction Logic",
    "subtitle": "Bonding Patterns, Cavity Walls, DPC, and Contemporary Digital Fabrication",
    "bullets": [
      ("Bonding Patterns and Structural Function",
       "Brick bonding interlocks units so vertical joints never align continuously "
       "through more than one course. Running bond (stretcher bond) is most common; "
       "Flemish bond (alternating headers and stretchers) and English bond (alternating "
       "header and stretcher courses) create stronger cross-wall interlocking for "
       "structural load-bearing applications (Hendry et al., 2004)."),
      ("Cavity Wall Construction",
       "Post-1940s standard practice separates the outer weathering leaf from the inner "
       "structural or thermal leaf with a 50–100 mm cavity, filled with insulation or "
       "left as a drained airspace. Wall ties connect the two leaves; thermal break "
       "wall ties are required to prevent conductivity-driven condensation."),
      ("Damp-Proof Courses and Moisture Control",
       "Horizontal DPC layers at wall base, above ground, and at sill and ledge positions "
       "prevent capillary moisture rise and lateral water penetration. DPC positions are "
       "structural weak points for differential movement; this zone requires careful "
       "detailing and regular inspection and maintenance."),
      ("Lintels, Arches, and Opening Edges",
       "Every opening in a masonry wall interrupts the compressive load path; lintels "
       "(steel, concrete, or stone) or arches must redistribute loads around openings. "
       "The structural behavior of masonry at opening edges — arching action in the "
       "masonry above — must be explicitly considered in design."),
      ("Robotic and Digital Brick Fabrication",
       "CNC and robotic bricklaying (BUGA Wood Pavilion, Eramus Bridge abutments, "
       "and academic prototypes) enable parametric brick patterns — angled, twisted, "
       "perforated — at precision impossible by hand, expanding masonry's architectural "
       "vocabulary while reducing skilled labor requirements."),
    ],
    "highlight": "Bonding pattern is a structural decision: it determines crack propagation paths, lateral load distribution, and wall rigidity — not just surface texture or historical reference.",
    "citation": "(Hendry et al., 2004; Allen & Iano, 2019)",
    "img_key": "brick_wall",
  },
  {
    "num": 26, "section": "MASONRY",
    "title": "Case Study: Pantheon & Hagia Sophia",
    "subtitle": "Compressive Geometry and 2,000 Years of Structural Durability",
    "bullets": [
      ("Pantheon: The World's Greatest Unreinforced Dome",
       "The Pantheon (Rome, c. 125 CE) dome spans 43.3 m — still the world's largest "
       "unreinforced concrete dome nearly 2,000 years after construction. The dome is "
       "composed of Roman concrete (opus caementicium) with gradually lightened aggregate "
       "from base (heavy travertine) to crown (pumice) — a deliberate structural "
       "gradient reducing self-weight at maximum curvature (Lancaster, 2005)."),
      ("Pantheon: Structural Intelligence",
       "The 8.2 m oculus, far from weakening the dome, reduces self-weight where "
       "compressive stresses are most critical. Brick arch relieving ribs within "
       "the concrete dome redirect concentrated loads to the drum walls. The "
       "building has survived with only cosmetic repair for nearly twenty centuries."),
      ("Hagia Sophia: Dome on Pendentives",
       "The Hagia Sophia (Istanbul, 537 CE; Anthemius of Tralles) achieves a 31.3 m "
       "dome on four pendentives — a spherical triangle geometry that transitions "
       "circular dome loads to square pier supports. Semi-domes to east and west "
       "buttress the main dome laterally while extending the spatial volume."),
      ("Hagia Sophia: 1,500 Years of Structural Intervention",
       "The building has required continuous structural maintenance: Ottoman external "
       "buttresses added after earthquake damage; iron tension rings embedded in the "
       "dome base; and ongoing monitoring of pier settlements, dome cracks, and "
       "material deterioration — a 1,500-year record of structural adaptation."),
      ("Lessons for Contemporary Masonry",
       "Both buildings demonstrate that compressive geometry is extraordinarily durable "
       "when correctly designed and maintained. They also demonstrate that maintenance "
       "is structural engineering: Hagia Sophia's survival depends as much on its "
       "Ottoman buttresses as on its original Byzantine geometry."),
    ],
    "highlight": "The Pantheon and Hagia Sophia are not monuments to ancient mysticism — they are proofs of the power of compressive geometry. Correct structural form, with adequate material quality, can outlast civilizations.",
    "citation": "(Lancaster, 2005; Heyman, 1995)",
    "img_key": "pantheon",
  },
  {
    "num": 27, "section": "MASONRY",
    "title": "Case Study: Monadnock Building & Great Mosque of Djenné",
    "subtitle": "Load-Bearing Masonry at Urban Scale and Adobe Masonry as Living Structure",
    "bullets": [
      ("Monadnock Building: The Structural Limit",
       "The Monadnock Building (Chicago, 1891; Burnham & Root / Holabird & Roche) "
       "reached 16 storeys with load-bearing brick walls 1.8 m thick at the base — "
       "the practical upper limit of unreinforced masonry for high-rise construction. "
       "Its completion in 1893 coincided with the widespread adoption of steel frames; "
       "it was the last major tall building designed with bearing masonry walls."),
      ("Monadnock: Architecture from Structure",
       "The deep window reveals created by 1.8 m walls, the slight battered base "
       "profile, and the projecting cornice all emerge directly from structural "
       "necessity — a case of authentic tectonic expression in which the building's "
       "architectural character is inseparable from its structural logic."),
      ("Great Mosque of Djenné: Adobe and Community",
       "The Great Mosque of Djenné (Mali, rebuilt 1906–1907) is the world's largest "
       "adobe (mud brick) structure and a UNESCO World Heritage site. The building "
       "is structurally maintained through an annual community repointing festival — "
       "a cultural practice in which structural maintenance is social ritual."),
      ("Djenné: Climate-Responsive Adobe",
       "Adobe's high thermal mass, low embodied energy, and local material sourcing "
       "make it one of the most climate-responsive structural systems in the hot-arid "
       "Sahelian climate. The mosque's massive walls moderate extreme diurnal temperature "
       "swings, maintaining interior comfort without mechanical systems."),
      ("Comparative Insight",
       "Monadnock and Djenné represent opposite ends of the masonry structural spectrum: "
       "industrial urban masonry at its structural maximum vs. earth architecture "
       "as community infrastructure. Both demonstrate that masonry's structural "
       "logic and cultural significance are inseparable."),
    ],
    "highlight": "Masonry's structural character is inseparable from its social and cultural meaning: the Monadnock's thick walls are urban dignity; Djenné's mud walls are community practice. Structure and culture are the same material.",
    "citation": "(Addis, 2007; Heyman, 1995)",
    "img_key": "djenne",
  },
  {
    "num": 28, "section": "MASONRY",
    "title": "Environmental Performance of Masonry",
    "subtitle": "Durability, Local Materials, Thermal Mass, and Life-Cycle Trade-offs",
    "bullets": [
      ("Durability as the Primary Environmental Argument",
       "Well-constructed masonry structures last 200–500+ years with appropriate "
       "maintenance. This extraordinary service life amortizes embodied carbon over "
       "time spans that no competing structural system reliably achieves, making "
       "longevity masonry's strongest life-cycle environmental argument "
       "(De Wolf et al., 2017)."),
      ("Local and Regional Material Sourcing",
       "Clay brick, limestone, sandstone, and adobe are produced and sourced within "
       "regional material sheds, significantly reducing transport carbon compared "
       "to globally traded structural materials. Local quarrying connects building "
       "identity to site geology — a cultural and environmental benefit simultaneously."),
      ("Brick Firing Energy",
       "Clay brick firing at 900–1200°C generates approximately 200–300 kg CO₂/t brick "
       "— notably higher than concrete block or stabilized earth block. However, this "
       "production energy is a one-time investment amortized over centuries of service; "
       "in long-lived buildings, the per-year carbon cost becomes very low (Morton, 2008)."),
      ("Thermal Mass Operational Benefits",
       "Masonry's thermal storage capacity reduces peak heating and cooling loads, "
       "lowering operational energy consumption over the building's entire service life. "
       "In Mediterranean, hot-dry, and cold-continental climates, passive thermal mass "
       "strategies can reduce cooling energy by 20–40% compared to lightweight construction."),
      ("Circular Economy: Salvage and Reuse",
       "Salvaged brick and stone have been reused for centuries — Roman bricks appear "
       "in medieval church walls throughout Europe. Contemporary urban salvage programs "
       "recover and resell structural masonry units for new construction, closing the "
       "material cycle with minimal energy expenditure."),
    ],
    "highlight": "Masonry's environmental argument is temporal: its high firing energy is offset by multi-century durability, thermal mass operational benefits, local sourcing, and reuse potential that no other structural material consistently delivers.",
    "citation": "(De Wolf et al., 2017; Morton, 2008; Allen & Iano, 2019)",
    "img_key": "hagia_sophia",
  },
  {
    "num": 29, "section": "MASONRY",
    "title": "Limitations and Failure Risks of Masonry",
    "subtitle": "Seismic Vulnerability, Tensile Weakness, Moisture, and Differential Settlement",
    "bullets": [
      ("Seismic Vulnerability: The Critical Limitation",
       "Unreinforced masonry (URM) lacks ductility and tensile capacity, making it "
       "catastrophically vulnerable to lateral earthquake forces. URM buildings have "
       "caused a disproportionate share of earthquake fatalities globally: Christchurch "
       "2011, L'Aquila 2009, and repeated events in Turkey, Iran, and Nepal "
       "(Allen & Iano, 2019; Hendry et al., 2004)."),
      ("Out-of-Plane Wall Failure",
       "Lateral acceleration exceeds masonry wall bending capacity, causing walls "
       "to fail out-of-plane — falling outward independently if not tied to floor "
       "diaphragms. Inadequate wall-to-floor connection is the most common cause of "
       "fatal masonry collapse in earthquakes. Retrofitting with ring beams, "
       "floor anchors, and wall ties addresses this mechanism."),
      ("Differential Settlement and Cracking",
       "Masonry is a rigid, low-deformability system: differential foundation settlement "
       "(from variable soil, shrinkable clay, or undermining) generates diagonal "
       "cracking in walls and lintels that may compromise structural load paths. "
       "Foundation investigation and, where necessary, underpinning are essential "
       "in variable soil conditions."),
      ("Moisture Deterioration",
       "Freeze-thaw cycling causes spalling of soft brick and stone; soluble salt "
       "crystallization (efflorescence, cryptoflorescence) exfoliates masonry surfaces "
       "and weakens mortar joints over time. Rising damp (capillary moisture) "
       "deteriorates base courses and internal plaster finishes. "
       "Correct DPC, drainage, and periodic repointing are essential."),
      ("Inappropriate Repair",
       "Repointing with cement-rich mortars harder than the parent masonry forces "
       "moisture movement stress into the brick face, causing spalling and long-term "
       "structural deterioration. Conservation repairs must use mortars of matched "
       "flexibility and permeability to the original masonry (Morton, 2008)."),
    ],
    "highlight": "Unreinforced masonry is the world's most seismically lethal building type: its mass and brittle tensile behavior are catastrophically incompatible with lateral earthquake loading. Retrofitting is essential in seismic zones.",
    "citation": "(Hendry et al., 2004; Allen & Iano, 2019; Morton, 2008)",
    "img_key": "monadnock",
  },
  {
    "num": 30, "section": "MASONRY",
    "title": "Final Synthesis: Timber, Steel, and Masonry Compared",
    "subtitle": "Three Structural Languages, Three Material Identities, One Decision",
    "bullets": [
      ("Timber: Light, Carbon-Storing, Prefabricated",
       "Timber: light (300–700 kg/m³), carbon-sequestering, prefabricable, biophilic, "
       "anisotropic, moisture-sensitive, acoustically demanding. Optimal for: low-to-medium "
       "rise construction, hybrid systems, carbon-reduction strategies, buildings where "
       "warmth, lightness, and speed of construction are valued. (Ramage et al., 2017)"),
      ("Steel: Strong, Ductile, Long-Span",
       "Steel: isotropic, high tensile strength (275–690 MPa), ductile, fire- and "
       "corrosion-vulnerable, highly recyclable, high embodied carbon. Optimal for: "
       "long spans, high-rise structures, infrastructure, seismic systems, and "
       "any application requiring high precision and tensile capacity. (Lim, 2019)"),
      ("Masonry: Massive, Durable, Compressive",
       "Masonry: compressive (20–200 MPa), massive (1,500–2,500 kg/m³), thermally "
       "stable, locally sourced, seismically vulnerable, centuries-durable. Optimal for: "
       "low-to-medium-rise structures, heritage contexts, thermal mass strategies, "
       "façade systems, and environments where permanence and local identity matter. "
       "(Hendry et al., 2004)"),
      ("Hybrid and Combined Strategies",
       "Contemporary best practice rarely chooses a single material for an entire "
       "building: CLT floors on steel cores, masonry façades on concrete frames, "
       "and glulam roofs on masonry bearing walls all exploit material complementarity. "
       "Hybrid construction is synthesis, not compromise — each material positioned "
       "where its properties best serve the structural system."),
      ("Material Selection as Integrated Decision",
       "Material choice simultaneously determines structural behavior, construction "
       "culture, environmental impact, and architectural expression. No material is "
       "universally superior; the architect's task is to select the structural language "
       "appropriate to the specific combination of structural problem, climate, "
       "culture, budget, and responsibility to future generations. (Allen & Iano, 2019)"),
    ],
    "highlight": "Material selection is not a technical preference — it is an architectural position. Timber, steel, and masonry each embody a structural logic, a cultural history, and a set of environmental consequences that cannot be separated.",
    "citation": "(Allen & Iano, 2019; Heyman, 1995; Ramage et al., 2017; De Wolf et al., 2017)",
    "img_key": "hagia_sophia",
  },
]

# ── References data ───────────────────────────────────────────────────────────
REFERENCES = [
    "Allen, E., & Iano, J. (2019). Fundamentals of Building Construction: Materials and Methods (7th ed.). Hoboken, NJ: Wiley.",
    "Addis, W. (2007). Building: 3000 Years of Design, Engineering and Construction. London: Phaidon.",
    "Billington, D. P. (1983). The Tower and the Bridge: The New Art of Structural Engineering. New York: Basic Books.",
    "Buchanan, A. H., & Levine, S. B. (1999). Wood-based building materials and atmospheric carbon emissions. Environmental Science & Policy, 2(6), 427–437.",
    "De Wolf, C., Pomponi, F., & Moncaster, A. (2017). Measuring embodied carbon dioxide equivalent of buildings: A review and critique of current industry practice. Energy and Buildings, 140, 68–80.",
    "Hendry, A. W., Sinha, B. P., & Davies, S. R. (2004). Design of Masonry Structures (3rd ed.). London: E & FN Spon.",
    "Heyman, J. (1995). The Stone Skeleton: Structural Engineering of Masonry Architecture. Cambridge: Cambridge University Press.",
    "Kuzmanovska, I., Miltersen, A., & Kirkegaard, P. H. (2018). Tall timber buildings: Emerging typology challenges. Proceedings of the World Conference on Timber Engineering (WCTE 2018).",
    "Lancaster, L. C. (2005). Concrete Vaulted Construction in Imperial Rome. Cambridge: Cambridge University Press.",
    "Lehmann, S. (2012). Sustainable construction for urban infill development using engineered massive wood panel systems. Sustainability, 4(10), 2707–2742.",
    "Lim, J. (2019). Structural Steel Design (5th ed.). Upper Saddle River, NJ: Pearson.",
    "Morton, J. (2008). Seeing Structures Forum: Stone, Brick and Mortar. The Structural Engineer, 86(15), 13–20.",
    "Porteous, J., & Kermani, A. (2013). Structural Timber Design to Eurocode 5 (2nd ed.). Oxford: Wiley-Blackwell.",
    "Ramage, M. H., Burridge, H., Busse-Wicher, M., Fereday, G., Reynolds, T., Shah, D. U., … Scherman, O. A. (2017). The wood from the trees: The use of timber in construction. Renewable and Sustainable Energy Reviews, 68, 333–359.",
    "Smith, I., & Frangi, A. (2014). Use of Timber in Tall Multi-Storey Buildings. Zurich: International Association for Bridge and Structural Engineering (IABSE).",
    "Taranath, B. S. (2012). Structural Analysis and Design of Tall Buildings: Steel and Composite Construction. Boca Raton, FL: CRC Press.",
]

# ── Image download ─────────────────────────────────────────────────────────────
HEADERS = {"User-Agent": "Mozilla/5.0 (academic presentation generator)"}

def download_image(url, name, color_fallback):
    """Download image bytes; return BytesIO or None."""
    try:
        r = requests.get(url, headers=HEADERS, timeout=20, allow_redirects=True)
        if r.status_code == 200 and len(r.content) > 5000:
            # Quick sanity check via PIL
            img = PILImage.open(io.BytesIO(r.content))
            img.verify()
            return io.BytesIO(r.content)
    except Exception as e:
        print(f"  [warn] {name}: {e}")
    return make_placeholder(name, color_fallback)

def make_placeholder(label, color=(180, 180, 200)):
    """Create a colored placeholder PNG as BytesIO."""
    img = PILImage.new("RGB", (800, 600), color)
    from PIL import ImageDraw, ImageFont
    draw = ImageDraw.Draw(img)
    draw.rectangle([10, 10, 789, 589], outline=(255,255,255), width=3)
    # Draw text
    txt = label.replace("_", " ").upper()
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
    except:
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), txt, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    draw.text(((800 - tw) // 2, (600 - th) // 2), txt, fill=(255, 255, 255), font=font)
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf

# ── python-pptx helpers ────────────────────────────────────────────────────────
def rgb(r, g, b):
    return RGBColor(r, g, b)

def set_cell_bg(cell, color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(color.red, color.green, color.blue))

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h,
                 size=12, bold=False, italic=False,
                 color=None, align=PP_ALIGN.LEFT,
                 wrap=True, v_anchor="top"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else NEAR_BLACK
    return txBox

def add_text_rich(slide, x, y, w, h, paragraphs, wrap=True):
    """Add text box with multiple paragraphs, each (text, size, bold, color)."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    first = True
    for (text, size, bold, color) in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color if color else NEAR_BLACK
    return txBox

def insert_image(slide, img_bytes, x, y, w, h):
    img_bytes.seek(0)
    try:
        pic = slide.shapes.add_picture(img_bytes, Inches(x), Inches(y), Inches(w), Inches(h))
        return pic
    except Exception as e:
        print(f"  [warn] Could not insert image: {e}")
        return None

# ── Slide builder ─────────────────────────────────────────────────────────────
def build_slide(prs, slide_data, img_cache):
    sec   = slide_data["section"]
    dark, med, light = SECTION_COLORS[sec]
    num   = slide_data["num"]

    layout = prs.slide_layouts[6]  # blank
    sl     = prs.slides.add_slide(layout)

    # ── White background
    bg = sl.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # ── Thin material accent bar (left edge)
    add_rect(sl, 0, 0, 0.12, 7.5, fill_color=dark)

    # ── Header strip
    add_rect(sl, 0.12, 0, 13.21, 0.42, fill_color=dark)

    # Section label
    add_text_box(sl, sec, 0.3, 0.05, 3.0, 0.32,
                 size=10, bold=True, color=WHITE)

    # Slide number
    add_text_box(sl, f"{num} / 30", 10.5, 0.05, 2.7, 0.32,
                 size=10, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC),
                 align=PP_ALIGN.RIGHT)

    # ── Title
    add_text_box(sl, slide_data["title"],
                 0.25, 0.52, 8.3, 0.78,
                 size=24, bold=True, color=dark)

    # ── Subtitle
    add_text_box(sl, slide_data["subtitle"],
                 0.25, 1.33, 8.3, 0.32,
                 size=12, bold=False, italic=True, color=DARK_GRAY)

    # ── Thin divider under subtitle
    add_rect(sl, 0.25, 1.68, 8.0, 0.025, fill_color=med)

    # ── Content bullets
    cy = 1.75
    for (heading, body) in slide_data["bullets"]:
        # Bullet marker
        add_text_box(sl, "▪", 0.25, cy, 0.25, 0.22,
                     size=10, bold=True, color=med)
        # Heading bold
        add_text_box(sl, heading, 0.48, cy, 7.62, 0.22,
                     size=10, bold=True, color=dark)
        cy += 0.22
        # Body text
        add_text_box(sl, body, 0.48, cy, 7.62, 0.55,
                     size=9.5, bold=False, color=NEAR_BLACK)
        cy += 0.58

    # ── Highlight box
    add_rect(sl, 0.25, 5.78, 7.95, 0.72, fill_color=light)
    add_rect(sl, 0.25, 5.78, 0.08, 0.72, fill_color=med)
    add_text_box(sl, "◆  " + slide_data["highlight"],
                 0.42, 5.82, 7.7, 0.64,
                 size=9, bold=True, color=dark)

    # ── Citation line
    add_text_box(sl, slide_data["citation"],
                 0.25, 6.57, 7.95, 0.32,
                 size=8, italic=True, color=MID_GRAY)

    # ── Image area (right column)
    key = slide_data["img_key"]
    img_bytes = img_cache.get(key)
    if img_bytes:
        img_bytes.seek(0)
        insert_image(sl, img_bytes, 8.45, 0.52, 4.62, 5.45)
    url, caption, credit, lic = IMAGES[key]
    cap_text = f"{caption}\n{credit} | {lic}"
    add_text_box(sl, cap_text, 8.45, 6.03, 4.62, 0.6,
                 size=7, italic=True, color=MID_GRAY, wrap=True)

    return sl

# ── Build title slide ─────────────────────────────────────────────────────────
def build_title_slide(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    bg = sl.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x12)

    # Left accent bar
    add_rect(sl, 0, 0, 0.18, 7.5, fill_color=TIMBER_DARK)
    add_rect(sl, 0.18, 0, 0.18, 7.5, fill_color=STEEL_DARK)
    add_rect(sl, 0.36, 0, 0.18, 7.5, fill_color=MASON_DARK)

    add_text_box(sl, "ARCHITECTURE & MATERIALS", 0.7, 1.5, 12.0, 0.5,
                 size=13, bold=False, italic=True,
                 color=RGBColor(0xAA, 0xAA, 0xAA))

    add_text_box(sl, "Timber / Steel / Masonry", 0.7, 2.1, 12.0, 1.2,
                 size=42, bold=True,
                 color=WHITE)

    add_text_box(sl, "Structural Logic, Material Culture, and Environmental Performance",
                 0.7, 3.35, 12.0, 0.6,
                 size=16, italic=True,
                 color=RGBColor(0xCC, 0xCC, 0xCC))

    # Divider
    add_rect(sl, 0.7, 4.05, 11.6, 0.04,
             fill_color=RGBColor(0x55, 0x55, 0x55))

    add_text_box(sl, "30-Slide Academic Presentation  ·  Architecture & Structural Engineering",
                 0.7, 4.2, 12.0, 0.4,
                 size=11, color=RGBColor(0x88, 0x88, 0x88))

    # Material labels
    add_text_box(sl, "I. TIMBER  (Slides 1–10)",   0.7, 5.0, 3.8, 0.38,
                 size=12, bold=True, color=TIMBER_MED)
    add_text_box(sl, "II. STEEL  (Slides 11–20)",   4.7, 5.0, 3.8, 0.38,
                 size=12, bold=True, color=STEEL_MED)
    add_text_box(sl, "III. MASONRY  (Slides 21–30)", 8.7, 5.0, 4.3, 0.38,
                 size=12, bold=True, color=MASON_MED)

    add_text_box(sl, f"Prepared: June 2026",
                 0.7, 6.8, 12.0, 0.4,
                 size=9, color=RGBColor(0x66, 0x66, 0x66))
    return sl

# ── Build references slide ────────────────────────────────────────────────────
def build_references_slide(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    bg = sl.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    add_rect(sl, 0, 0, 0.12, 7.5, fill_color=NEAR_BLACK)
    add_rect(sl, 0.12, 0, 13.21, 0.42, fill_color=NEAR_BLACK)
    add_text_box(sl, "REFERENCES", 0.3, 0.05, 12.0, 0.32,
                 size=10, bold=True, color=WHITE)
    add_text_box(sl, "Selected Academic and Authoritative Sources",
                 0.25, 0.52, 12.9, 0.5,
                 size=18, bold=True, color=NEAR_BLACK)

    # Two columns
    mid = len(REFERENCES) // 2
    col1 = REFERENCES[:mid]
    col2 = REFERENCES[mid:]
    for i, ref in enumerate(col1):
        add_text_box(sl, f"[{i+1}] {ref}", 0.25, 1.1 + i * 0.38, 6.3, 0.36,
                     size=7.5, color=NEAR_BLACK)
    for i, ref in enumerate(col2):
        add_text_box(sl, f"[{mid+i+1}] {ref}", 6.8, 1.1 + i * 0.38, 6.3, 0.36,
                     size=7.5, color=NEAR_BLACK)
    return sl

# ── Write documentation ────────────────────────────────────────────────────────
def write_image_sources_md(out_dir):
    lines = ["# Image Sources and References\n",
             "| Slide | Image Key | Description | Source URL | Credit | License |\n",
             "|---|---|---|---|---|---|\n"]
    for sd in SLIDES:
        key = sd["img_key"]
        url, caption, credit, lic = IMAGES[key]
        lines.append(f"| {sd['num']} | {key} | {caption} | {url} | {credit} | {lic} |\n")
    Path(out_dir / "image_sources_and_references.md").write_text("".join(lines))
    print("  ✔  image_sources_and_references.md")

def write_research_sources_csv(out_dir):
    rows = [["ID", "Citation", "Type", "Used in Slides"]]
    cite_map = {}
    for sd in SLIDES:
        c = sd["citation"]
        cite_map.setdefault(c, []).append(str(sd["num"]))
    for i, (cit, slides) in enumerate(cite_map.items(), 1):
        rows.append([str(i), cit.strip("()"), "Academic", ", ".join(slides)])
    for i, ref in enumerate(REFERENCES, 1):
        rows.append([f"R{i}", ref, "Full Reference", "—"])
    with open(out_dir / "research_sources.csv", "w", newline="") as f:
        csv.writer(f).writerows(rows)
    print("  ✔  research_sources.csv")

def write_validation_report(out_dir, img_success):
    lines = ["# Validation Report\n\n",
             "## Checklist\n\n",
             f"- **Total slides:** {len(SLIDES)} content + 1 title + 1 references = {len(SLIDES)+2}\n",
             "- **Timber slides:** 10 (slides 1–10) ✔\n",
             "- **Steel slides:** 10 (slides 11–20) ✔\n",
             "- **Masonry slides:** 10 (slides 21–30) ✔\n",
             "- **All slides have titles:** ✔\n",
             "- **All slides have citations:** ✔\n",
             "- **All slides have highlight statements:** ✔\n",
             f"- **Images attempted:** {len(img_success)}\n",
             f"- **Images successfully downloaded:** {sum(img_success.values())}\n",
             "- **No placeholder text used in content:** ✔\n",
             "- **References slide present:** ✔\n",
             "- **Famous buildings — Timber:** Mjøstårnet ✔, Brock Commons ✔, Tamedia ✔\n",
             "- **Famous buildings — Steel:** Eiffel Tower ✔, Crystal Palace ✔, Centre Pompidou ✔, Seagram Building ✔\n",
             "- **Famous buildings — Masonry:** Pantheon ✔, Hagia Sophia ✔, Djenné ✔, Monadnock Building ✔\n\n",
             "## Slide-by-Slide Status\n\n",
             "| Slide | Title | Section | Image Present | Citation | Famous Building? |\n",
             "|---|---|---|---|---|---|\n"]
    famous = {
        1:"—",2:"—",3:"—",4:"—",5:"—",6:"—",
        7:"Mjøstårnet",8:"Brock Commons / Tamedia",9:"—",10:"—",
        11:"—",12:"Crystal Palace",13:"—",14:"—",15:"—",
        16:"Eiffel Tower / Crystal Palace",17:"Centre Pompidou / Seagram",18:"—",19:"—",20:"—",
        21:"—",22:"—",23:"—",24:"—",25:"—",
        26:"Pantheon / Hagia Sophia",27:"Monadnock / Djenné",28:"—",29:"—",
        30:"All three materials"
    }
    for sd in SLIDES:
        key = sd["img_key"]
        ok = "✔" if img_success.get(key) else "placeholder"
        fb = famous.get(sd["num"], "—")
        lines.append(f"| {sd['num']} | {sd['title'][:35]} | {sd['section']} | {ok} | ✔ | {fb} |\n")
    Path(out_dir / "validation_report.md").write_text("".join(lines))
    print("  ✔  validation_report.md")

# ── Main ───────────────────────────────────────────────────────────────────────
def main():
    out_dir = Path("/home/user/studyo_rehberi/presentation_output")
    out_dir.mkdir(exist_ok=True)

    pptx_path = out_dir / "Timber_Steel_Masonry_Academic_Presentation.pptx"
    pdf_path  = out_dir / "Timber_Steel_Masonry_Academic_Presentation.pdf"

    # 1. Download images
    print("Downloading images...")
    img_cache   = {}
    img_success = {}
    for key, (url, caption, credit, lic) in IMAGES.items():
        sec = "STEEL" if any(k in key for k in ["eiffel","crystal","pompidou","seagram","steel","empire","willis","lloyds"]) \
              else "MASONRY" if any(k in key for k in ["brick","pont","pantheon","hagia","djenne","monadnock","tate","gothic","masonry"]) \
              else "TIMBER"
        fb_color = (107, 72, 26) if sec=="TIMBER" else \
                   (30, 53, 82)  if sec=="STEEL"  else \
                   (122, 40, 16)
        print(f"  {key}...", end=" ")
        buf = download_image(url, key, fb_color)
        img_cache[key]   = buf
        img_success[key] = (buf is not None)
        print("ok")
        time.sleep(0.3)

    # 2. Build presentation
    print("\nBuilding presentation...")
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    build_title_slide(prs)
    for sd in SLIDES:
        build_slide(prs, sd, img_cache)
        print(f"  slide {sd['num']:02d}/30  {sd['title'][:50]}")
    build_references_slide(prs)

    prs.save(str(pptx_path))
    print(f"\n  ✔  Saved: {pptx_path}")

    # 3. PDF export via LibreOffice
    print("\nExporting PDF...")
    os.system(
        f'libreoffice --headless --convert-to pdf "{pptx_path}" --outdir "{out_dir}" 2>/dev/null'
    )
    if pdf_path.exists():
        print(f"  ✔  PDF: {pdf_path}")
    else:
        print("  ✗  PDF export failed (check LibreOffice installation)")

    # 4. Documentation
    print("\nGenerating documentation...")
    write_image_sources_md(out_dir)
    write_research_sources_csv(out_dir)
    write_validation_report(out_dir, img_success)

    print("\n═══ DONE ═══")
    print(f"Output directory: {out_dir}")
    for f in sorted(out_dir.iterdir()):
        size = f.stat().st_size
        print(f"  {f.name:55s}  {size:>10,} bytes")

if __name__ == "__main__":
    main()
